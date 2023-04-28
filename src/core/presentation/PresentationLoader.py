import collections
import io
from collections import abc
import pptx
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml import parse_xml
from pptx import Presentation
from pptx.presentation import Presentation as PresentationClass
from pptx.slide import SlideLayout
from src.core.presentation.PresentationType import PresentationType
import src.core.presentation.PresentationPlaceholderCollections as PLC
from src.core.presentation.data.LoadedPresentation import LoadedPresentation
from src.core.presentation.data.PreLayout import PreLayout
from src.data.Fields import Fields
from src.translations.Translator import _

# Lookup names for the layouts
NAME_MASTER_REAL = "Präsentation"
NAME_MASTER_LOOKUP = "Spicker"

NAME_LAYOUT_LOOKUP_NORMAL = "Spicker_Preis"


# Loads the required masters from a presentation root
def __get_masters_from_presentation(root: PresentationClass):
    # Placeholder for both picks
    real = None
    lookup = None

    # Iterates over every slide-master to search for the required ones
    for sld_master in root.slide_masters:
        # Loads the real name
        theme = sld_master.part.part_related_by(RT.THEME)
        name = parse_xml(theme.blob).get("name")

        # Checks the name
        if name == NAME_MASTER_LOOKUP:
            lookup = sld_master
        elif name == NAME_MASTER_REAL:
            real = sld_master

    # Raises an error if one of the masters wasn't found
    if real is None or lookup is None:
        raise ValueError("master.notfound", (PresentationType.REAL if real is None else PresentationType.LOOKUP))

    return real, lookup


# Takes in a layout and a list of placeholders.
# Searches the required placeholders and returns them as a dict
# raises ValueError if a placeholder wasn't found or if the placeholder was of an invalid type
def __get_placeholders(layout: pptx.slide.SlideLayout, placeholder_group: [str]):
    # Holds all placeholders
    plc_dict = {}

    # Iterates over every required placeholder and then searches it
    for plc_text, plc_types in placeholder_group:
        for plc in layout.placeholders:
            # Gets the placeholder text
            txt: str = plc.text.strip().lower()

            # Ensures the texts match up
            if txt == plc_text.lower():

                # Ensures that the placeholder is of a valid type
                if plc.placeholder_format.type.real not in plc_types:
                    raise ValueError(_(
                        "The placeholder '{placeholder}' from the layout '{layout}' has an invalid format.").format(
                        placeholder=plc_text, layout=layout.name))

                # Appends the placeholder
                plc_dict[plc_text] = plc
                break

        # Ensures the placeholder is found
        if plc_text not in plc_dict:
            raise ValueError(
                _("The lookup-placeholder '{placeholder}' couldn't be found").format(placeholder=plc_text))

    return plc_dict


def __get_layout_with_name(master: pptx.slide.SlideMaster, name: str):
    # Iterates over all required lookup-layouts
    for layout in master.slide_layouts:

        # Checks the name
        if layout.name == name:
            return layout

    raise ValueError(_("The master-layout '{layout}' couldn't be found").format(layout=name))


# Loads the lookup-master slide
def __load_lookup_master(master: pptx.slide.SlideMaster):
    # Gets the layout
    layout = __get_layout_with_name(master, NAME_LAYOUT_LOOKUP_NORMAL)

    # Gets the placeholders
    placeholders = __get_placeholders(layout,
                                      [PLC.MEMBER_1, PLC.MEMBER_2, PLC.MEMBER_3, PLC.SCHOOL, PLC.NOTES,
                                       PLC.PROJECT_TITLE,
                                       PLC.LOOKUP_PRICE_1, PLC.LOOKUP_PRICE_2])

    return PreLayout(layout, placeholders)


def __load_real_master(master: pptx.slide.SlideMaster):
    def __load_presentation(fld: Fields):
        layout = __get_layout_with_name(master, fld.value[0])

        # Gets all those placeholders
        placeholders = __get_placeholders(layout, [
            # Members
            PLC.MEMBER_1, PLC.MEMBER_2, PLC.MEMBER_3,

            # Other stuff
            PLC.SCHOOL, PLC.PROJECT_TITLE, PLC.PROJECT_IMAGE,

            # Prices
            PLC.PROJECT_PRICES_1_FLYIN_TEXT, PLC.PROJECT_PRICES_1_FLYIN_IMAGE, PLC.PROJECT_PRICES_2_FLYIN_TEXT,
            PLC.PROJECT_PRICES_2_FLYIN_IMAGE, PLC.PROJECT_PRICES_1_ROTATING_TEXT, PLC.PROJECT_PRICES_1_ROTATING_IMAGE,
            PLC.PROJECT_PRICES_2_ROTATING_TEXT, PLC.PROJECT_PRICES_2_ROTATING_IMAGE
        ])

        return PreLayout(layout, placeholders)

    pres_dict: {Fields: PreLayout} = {}

    # Loads in all presentation
    for field in Fields:
        pres_dict[field] = __load_presentation(field)

    return pres_dict


# Takes in a path to the presentation and tries to load it
# Returns a tuple consisting of the following:
# (
#   {Fields: PreLayout}, # The layout for every field
#   PreLayout,           # The layout the lookup field
#   Presentation         # Well, the presentation
# )
#
# raises a ValueError if anything went wrong with loading the presentation or if anything is wrong with it's fields
#

def load_presentation(path: str):
    # Loads in the presentation
    root: PresentationClass = Presentation(path)

    # Gets the masters
    master_real, master_lookup = __get_masters_from_presentation(root)

    return LoadedPresentation(
        root,
        __load_real_master(master_real),
        __load_lookup_master(master_lookup),
    )

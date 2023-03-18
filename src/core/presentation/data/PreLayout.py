import pptx.shapes.placeholder
import pptx.slide


class PreLayout:

    def __init__(self, layout: pptx.slide.SlideLayout, placeholders: [pptx.shapes.placeholder.LayoutPlaceholder]):
        self.layout = layout
        self.placeholders = placeholders

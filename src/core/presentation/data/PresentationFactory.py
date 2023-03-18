import pptx.presentation
from pptx.dml.color import RGBColor

from src.core.presentation.data.PreLayout import PreLayout
from src.core.presentation.data.SlideFactory import SlideFactory
from src.data.Fields import Fields
from src.data.Price import DEFAUlT_PRICE, NormalPriceTypes, PriceAnimation
from src.data.Project import Project
from pptx.presentation import Presentation
import src.core.presentation.PresentationPlaceholderCollections as PLC
from src.data.ProjectMember import ProjectMember
from PIL import Image


class PresentationFactory:

    def __init__(self, presentation: Presentation, field_layouts: {Fields: PreLayout}, lookup_layout: PreLayout,
                 price_images: {NormalPriceTypes: Image}, special_price_image: Image, project_images: {str: Image}):
        self.field_layouts = field_layouts
        self.lookup_layout = lookup_layout
        self.presentation = presentation
        self.price_images = price_images
        self.special_price_image = special_price_image
        self.project_images = project_images

    # Takes in a price-animation and returns the four placeholders (price 1 (image,text), price 2 (image, text))
    # for that animation
    def __get_price_placeholders_by_animation(self, animation: PriceAnimation):
        if animation == PriceAnimation.FLYIN:
            return (
                PLC.PROJECT_PRICES_1_FLYIN_IMAGE,
                PLC.PROJECT_PRICES_1_FLYIN_TEXT,
                PLC.PROJECT_PRICES_2_FLYIN_IMAGE,
                PLC.PROJECT_PRICES_2_FLYIN_TEXT
            )
        else:
            return (
                PLC.PROJECT_PRICES_1_ROTATING_IMAGE,
                PLC.PROJECT_PRICES_1_ROTATING_TEXT,
                PLC.PROJECT_PRICES_2_ROTATING_IMAGE,
                PLC.PROJECT_PRICES_2_ROTATING_TEXT
            )

    # Takes in a project and build the project-slide for that project
    # If with_prices is enabled the prices for that project will also be added
    # and based on the animation-field will the correct animation be selected
    def create_project_slide(self, project: Project, with_prices: bool = False,
                             animation: PriceAnimation = PriceAnimation.FLYIN):
        # Adds the new slide
        factory = SlideFactory(self.presentation, self.field_layouts[project.field])

        # Gets the stand-id
        stand_id = project.get_raw_stand_number()

        # Checks if the project-image exists
        if stand_id in self.project_images:
            # Adds the image
            factory.add_image(PLC.PROJECT_IMAGE, self.project_images[stand_id])

        # Adds common information
        factory.populate(PLC.PROJECT_TITLE, project.name)
        factory.populate(PLC.SCHOOL, project.location)

        # Member references
        member_fields = [PLC.MEMBER_1, PLC.MEMBER_2, PLC.MEMBER_3]

        # Iterates over all project members to populate the slide with them
        for mid in range(len(project.members)):
            # Gets the member
            mem: ProjectMember = project.members[mid]
            factory.populate(member_fields[mid], mem.get_short_description())

        # Checks if prices should be loaded
        if with_prices and project.price is not DEFAUlT_PRICE:

            # Gets the price-placeholders
            plc_f_img, plc_f_text, plc_s_img, plc_s_text = self.__get_price_placeholders_by_animation(animation)

            # Adds the normal price (If one exists)
            if project.price.normal_price is not None:
                factory.add_image(plc_f_img, self.price_images[project.price.normal_price])

            # Adds the special price
            if project.price.has_special_price:
                # Gets the fields
                img_field = plc_f_img if project.price.normal_price is None else plc_s_img
                txt_field = plc_f_text if project.price.normal_price is None else plc_s_text

                factory.add_image(img_field, self.special_price_image)
                factory.populate(txt_field, project.special_price_name)

    # Takes in a project and build the lookup-slide for that project
    # If with_prices is enabled the prices for that project will also be added
    def create_lookup_slide(self, project: Project, with_prices: bool = False):
        # Adds the new slide
        factory = SlideFactory(self.presentation, self.lookup_layout)

        # Title
        factory.populate(PLC.PROJECT_TITLE, project.name)

        # Member references
        member_fields = [PLC.MEMBER_1, PLC.MEMBER_2, PLC.MEMBER_3]

        # Iterates over all project members to populate the slide with them
        for mid in range(len(project.members)):
            # Gets the member
            mem: ProjectMember = project.members[mid]

            # Gets the paragraph to add text to
            paragraph = factory.get_paragraph(member_fields[mid])

            # Adds firstname
            fn = paragraph.add_run()
            fn.text = mem.first_name
            fn.font.bold = True

            # Adds rest
            other = paragraph.add_run()
            other.text = " {} ({})".format(mem.last_name, mem.age)

        # School name
        factory.populate(PLC.SCHOOL, project.location)

        # Checks if prices should be loaded
        if with_prices and project.price is not DEFAUlT_PRICE:
            # Adds the nodes
            factory.populate(PLC.NOTES, project.price.get_name("\n+\n") + "\n" + project.special_price_name)

            # Adds the normal price (If one exists)
            if project.price.normal_price is not None:
                factory.add_image(PLC.LOOKUP_PRICE_1, self.price_images[project.price.normal_price])

            # Adds the special price
            if project.price.has_special_price:
                factory.add_image(PLC.LOOKUP_PRICE_2 if project.price.normal_price is None else PLC.LOOKUP_PRICE_2,
                                  self.special_price_image)

import pptx.presentation
from PIL import Image

from src.core.presentation.data.PreLayout import PreLayout
from pptx.slide import Slide
from pptx.shapes.placeholder import SlidePlaceholder
import io

class SlideFactory:

    def __init__(self, root: pptx.presentation.Presentation, base: PreLayout):
        self.base = base
        self.created: Slide = root.slides.add_slide(base.layout)

    # Gets the paragraph of a given placeholder to add additional formatted texts
    def get_paragraph(self, placeholder: (str, [int])):
        # Gets the placeholder
        plc = self.base.placeholders[placeholder[0]]

        # Gets the textbox ref
        ref: SlidePlaceholder = self.created.placeholders[plc.placeholder_format.idx]

        return ref.text_frame.paragraphs[0]

    # Populates a given placeholder with the given data (String)
    def populate(self, placeholder: (str, [int]), data: str):
        # Gets the placeholder
        plc = self.base.placeholders[placeholder[0]]

        # Gets the textbox ref
        ref: SlidePlaceholder = self.created.placeholders[plc.placeholder_format.idx]

        # Populates it with data
        ref.text = data

    def add_image(self, placeholder: (str, [int]), image: Image):
        # Gets the placeholder
        plc = self.base.placeholders[placeholder[0]]

        # Gets the textbox ref
        ref: SlidePlaceholder = self.created.placeholders[plc.placeholder_format.idx]

        # Converts the image to a file like object for the presentation to load
        b = io.BytesIO()
        image.save(b, "png")
        b.seek(0)

        ref.insert_picture(b)

